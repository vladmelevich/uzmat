#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Скрипт для создания презентации PowerPoint о проекте Uzmat
Требует установки: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

def create_presentation():
    """Создает премиальную презентацию о проекте Uzmat"""
    
    # Создаем презентацию
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Цветовая схема (премиальная темная тема)
    PRIMARY_COLOR = RGBColor(0, 150, 136)  # Зеленый (как в проекте)
    DARK_BG = RGBColor(15, 23, 42)  # Темно-синий фон
    LIGHT_TEXT = RGBColor(255, 255, 255)  # Белый текст
    ACCENT_COLOR = RGBColor(34, 197, 94)  # Яркий зеленый акцент
    
    # ========== СЛАЙД 1: Титульный ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой слайд
    
    # Фон
    background = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.color.rgb = DARK_BG
    
    # Заголовок
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "UZMAT"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Подзаголовок
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Премиальный маркетплейс спецтехники"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = LIGHT_TEXT
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Дата
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime("%d %B %Y")
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(16)
    date_para.font.color.rgb = RGBColor(150, 150, 150)
    date_para.alignment = PP_ALIGN.CENTER
    
    # ========== СЛАЙД 2: О проекте ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Фон
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = DARK_BG
    bg2.line.color.rgb = DARK_BG
    
    # Заголовок слайда
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title2_frame = title2.text_frame
    title2_frame.text = "О ПРОЕКТЕ"
    title2_frame.paragraphs[0].font.size = Pt(44)
    title2_frame.paragraphs[0].font.bold = True
    title2_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Описание
    desc_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    
    description = """Uzmat — это премиальная платформа для поиска и размещения объявлений о спецтехнике в Центральной Азии.

🎯 ОСНОВНАЯ ЦЕЛЬ
Единая точка доступа для поиска надежных подрядчиков и управления парком техники

✨ ОСОБЕННОСТИ
• Современный премиальный дизайн
• Полная адаптивность (мобильные, планшеты, десктоп)
• Плавные анимации и интуитивный интерфейс
• Многоязычная поддержка (KZ, UZ, RU, BY)
• Безопасная система авторизации

🌍 ГЕОГРАФИЯ
Казахстан • Узбекистан • Россия"""
    
    desc_frame.text = description
    for para in desc_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(12)
    
    # ========== СЛАЙД 3: Основные функции ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = DARK_BG
    bg3.line.color.rgb = DARK_BG
    
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title3_frame = title3.text_frame
    title3_frame.text = "ОСНОВНЫЕ ФУНКЦИИ"
    title3_frame.paragraphs[0].font.size = Pt(44)
    title3_frame.paragraphs[0].font.bold = True
    title3_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Функции в две колонки
    functions_left = [
        "🏗️ Продажа спецтехники",
        "🚛 Аренда спецтехники",
        "🔧 Запчасти",
        "🛠️ Ремонт и услуги"
    ]
    
    functions_right = [
        "💬 Система чатов",
        "⭐ Избранное",
        "👤 Профиль пользователя",
        "🔍 Умный поиск и фильтры"
    ]
    
    # Левая колонка
    left_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.text = "\n".join(functions_left)
    for para in left_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(16)
        para.font.bold = True
    
    # Правая колонка
    right_box = slide3.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.text = "\n".join(functions_right)
    for para in right_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(16)
        para.font.bold = True
    
    # ========== СЛАЙД 4: Главная страница ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = DARK_BG
    bg4.line.color.rgb = DARK_BG
    
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title4_frame = title4.text_frame
    title4_frame.text = "ГЛАВНАЯ СТРАНИЦА"
    title4_frame.paragraphs[0].font.size = Pt(44)
    title4_frame.paragraphs[0].font.bold = True
    title4_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content4_frame = content4.text_frame
    content4_frame.word_wrap = True
    
    main_page_text = """📍 URL: / (главная страница)

🎨 ДИЗАЙН
• Премиальный темный интерфейс
• Hero-секция с поиском
• Быстрые кнопки популярных запросов

📋 СЕКЦИИ
1. Горячие предложения
   - Последние объявления о продаже и аренде
   - Адаптивная сетка карточек
   
2. Популярные предложения
   - Фильтруемые объявления
   - Динамическая загрузка через AJAX
   
3. Категории
   • Продажа спецтехники
   • Аренда спецтехники
   • Заявки
   • Запчасти
   • Ремонт

🔍 ФУНКЦИОНАЛ
• Поиск в реальном времени
• Фильтры по стране, городу, типу, марке, цене
• AJAX-фильтрация без перезагрузки страницы"""
    
    content4_frame.text = main_page_text
    for para in content4_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 5: Каталог ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = DARK_BG
    bg5.line.color.rgb = DARK_BG
    
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title5_frame = title5.text_frame
    title5_frame.text = "КАТАЛОГ ОБЪЯВЛЕНИЙ"
    title5_frame.paragraphs[0].font.size = Pt(44)
    title5_frame.paragraphs[0].font.bold = True
    title5_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content5_frame = content5.text_frame
    content5_frame.word_wrap = True
    
    catalog_text = """📍 URL: /catalog/

🎯 НАЗНАЧЕНИЕ
Централизованный каталог всех объявлений о продаже и аренде спецтехники

📊 РЕЖИМЫ ОТОБРАЖЕНИЯ
• Все объявления (продажа + аренда) - по умолчанию
• Только продажа - /catalog/?ad_type=sale
• Только аренда - /catalog/?ad_type=rent

🔍 ФИЛЬТРЫ
✓ Страна (Казахстан, Узбекистан, Россия)
✓ Город (динамический список)
✓ Тип техники
✓ Марка
✓ Цена (от/до)

⚡ ТЕХНОЛОГИИ
• AJAX-фильтрация без перезагрузки
• Пагинация (12 объявлений на страницу)
• Динамическое обновление заголовков
• Сохранение фильтров в URL"""
    
    content5_frame.text = catalog_text
    for para in content5_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 6: Запчасти ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = DARK_BG
    bg6.line.color.rgb = DARK_BG
    
    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title6_frame = title6.text_frame
    title6_frame.text = "ЗАПЧАСТИ"
    title6_frame.paragraphs[0].font.size = Pt(44)
    title6_frame.paragraphs[0].font.bold = True
    title6_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content6_frame = content6.text_frame
    content6_frame.word_wrap = True
    
    parts_text = """📍 URL: /parts/

🎯 НАЗНАЧЕНИЕ
Специализированный раздел для поиска и размещения объявлений о запчастях для спецтехники

🔍 ПОИСК
• Поиск по названию запчасти
• Быстрые кнопки популярных запросов:
  - Фильтр
  - Масло
  - Гидроцилиндр
  - Тормозные колодки
  - Ремень
  - Шланг

📋 ФИЛЬТРЫ
✓ Страна и город
✓ Тип техники
✓ Марка
✓ Цена

📊 ОТОБРАЖЕНИЕ
• Сетка карточек объявлений
• Пагинация
• Динамическая фильтрация через AJAX"""
    
    content6_frame.text = parts_text
    for para in content6_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 7: Ремонт ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg7.fill.solid()
    bg7.fill.fore_color.rgb = DARK_BG
    bg7.line.color.rgb = DARK_BG
    
    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title7_frame = title7.text_frame
    title7_frame.text = "РЕМОНТ И УСЛУГИ"
    title7_frame.paragraphs[0].font.size = Pt(44)
    title7_frame.paragraphs[0].font.bold = True
    title7_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content7 = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content7_frame = content7.text_frame
    content7_frame.word_wrap = True
    
    repair_text = """📍 URL: /logistics/

🎯 НАЗНАЧЕНИЕ
Раздел для поиска и размещения услуг по ремонту и обслуживанию спецтехники

📋 КОНТЕНТ
• Отображает все объявления типа "Услуги"
• Специализированные сервисы по ремонту
• Профессиональные мастера и сервисные центры

🔍 ФИЛЬТРЫ
✓ Страна и город
✓ Тип техники
✓ Марка
✓ Цена

⚡ ОСОБЕННОСТИ
• Динамическая фильтрация
• Пагинация результатов
• Прямой контакт с исполнителями"""
    
    content7_frame.text = repair_text
    for para in content7_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 8: Профиль и авторизация ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg8.fill.solid()
    bg8.fill.fore_color.rgb = DARK_BG
    bg8.line.color.rgb = DARK_BG
    
    title8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title8_frame = title8.text_frame
    title8_frame.text = "ПРОФИЛЬ И АВТОРИЗАЦИЯ"
    title8_frame.paragraphs[0].font.size = Pt(44)
    title8_frame.paragraphs[0].font.bold = True
    title8_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content8 = slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content8_frame = content8.text_frame
    content8_frame.word_wrap = True
    
    auth_text = """👤 РЕГИСТРАЦИЯ
• Выбор типа: физическое лицо или компания
• Быстрая регистрация с минимальными данными
• Расширенные возможности для компаний

🔐 АВТОРИЗАЦИЯ
• Безопасный вход по email и паролю
• Двухфакторная аутентификация (опционально)
• Восстановление пароля

📱 ПРОФИЛЬ ПОЛЬЗОВАТЕЛЯ
📍 URL: /profile/

ФУНКЦИИ:
✓ Просмотр своих объявлений
✓ Управление объявлениями
✓ Редактирование профиля
✓ Настройки аккаунта
✓ История активности"""
    
    content8_frame.text = auth_text
    for para in content8_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 9: Создание объявления ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg9.fill.solid()
    bg9.fill.fore_color.rgb = DARK_BG
    bg9.line.color.rgb = DARK_BG
    
    title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title9_frame = title9.text_frame
    title9_frame.text = "СОЗДАНИЕ ОБЪЯВЛЕНИЯ"
    title9_frame.paragraphs[0].font.size = Pt(44)
    title9_frame.paragraphs[0].font.bold = True
    title9_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content9 = slide9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content9_frame = content9.text_frame
    content9_frame.word_wrap = True
    
    create_text = """📍 URL: /create/

📝 ТИПЫ ОБЪЯВЛЕНИЙ
• Продажа спецтехники
• Аренда спецтехники
• Запчасти
• Услуги (ремонт)

📋 ФОРМА ЗАПОЛНЕНИЯ
✓ Заголовок и описание
✓ Тип техники и марка
✓ Год выпуска и состояние
✓ Цена и валюта
✓ Местоположение (страна, город)
✓ Контактная информация
✓ Загрузка фотографий (множественная)

⚡ ОСОБЕННОСТИ
• Валидация данных
• Автоматическое создание slug
• Предпросмотр перед публикацией
• Сохранение черновиков"""
    
    content9_frame.text = create_text
    for para in content9_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 10: Детальная страница ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg10.fill.solid()
    bg10.fill.fore_color.rgb = DARK_BG
    bg10.line.color.rgb = DARK_BG
    
    title10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title10_frame = title10.text_frame
    title10_frame.text = "ДЕТАЛЬНАЯ СТРАНИЦА ОБЪЯВЛЕНИЯ"
    title10_frame.paragraphs[0].font.size = Pt(44)
    title10_frame.paragraphs[0].font.bold = True
    title10_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content10 = slide10.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content10_frame = content10.text_frame
    content10_frame.word_wrap = True
    
    detail_text = """📍 URL: /ad/<slug>/

📸 ГАЛЕРЕЯ
• Множественные фотографии
• Увеличение по клику
• Навигация между фото

📋 ИНФОРМАЦИЯ
✓ Полное описание
✓ Характеристики техники
✓ Цена и условия
✓ Местоположение
✓ Контактные данные

⭐ ФУНКЦИИ
• Добавление в избранное
• Поделиться объявлением
• Связаться с продавцом
• Просмотр других объявлений пользователя

📊 СТАТИСТИКА
• Счетчик просмотров
• Дата публикации
• Дата обновления"""
    
    content10_frame.text = detail_text
    for para in content10_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 11: Чаты и заявки ==========
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg11 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg11.fill.solid()
    bg11.fill.fore_color.rgb = DARK_BG
    bg11.line.color.rgb = DARK_BG
    
    title11 = slide11.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title11_frame = title11.text_frame
    title11_frame.text = "ЧАТЫ И ЗАЯВКИ"
    title11_frame.paragraphs[0].font.size = Pt(44)
    title11_frame.paragraphs[0].font.bold = True
    title11_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content11 = slide11.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content11_frame = content11.text_frame
    content11_frame.word_wrap = True
    
    chats_text = """📍 URL: /chats/

💬 ФУНКЦИОНАЛ
• Прямая связь между покупателем и продавцом
• Обмен сообщениями в реальном времени
• История переписки

📋 ЗАЯВКИ
• Создание заявок на аренду
• Управление заявками
• Статусы заявок (новая, в работе, выполнена)

🔔 УВЕДОМЛЕНИЯ
• Новые сообщения
• Обновления заявок
• Ответы на объявления

⚡ ОСОБЕННОСТИ
• Удобный интерфейс чата
• Поиск по перепискам
• Фильтрация заявок"""
    
    content11_frame.text = chats_text
    for para in content11_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 12: Технологии ==========
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg12.fill.solid()
    bg12.fill.fore_color.rgb = DARK_BG
    bg12.line.color.rgb = DARK_BG
    
    title12 = slide12.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title12_frame = title12.text_frame
    title12_frame.text = "ТЕХНОЛОГИИ И ИНСТРУМЕНТЫ"
    title12_frame.paragraphs[0].font.size = Pt(44)
    title12_frame.paragraphs[0].font.bold = True
    title12_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content12 = slide12.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content12_frame = content12.text_frame
    content12_frame.word_wrap = True
    
    tech_text = """🖥️ BACKEND
• Django 4.2.7
• Python 3.12
• SQLite (разработка)
• Django REST Framework (готово к интеграции)

🎨 FRONTEND
• HTML5, CSS3, JavaScript
• Адаптивный дизайн
• AJAX для динамической загрузки
• Плавные анимации

🗄️ БАЗА ДАННЫХ
• Модели: User, Advertisement, Category, Favorite
• Связи и индексы для оптимизации
• Миграции Django

🔒 БЕЗОПАСНОСТЬ
• Аутентификация Django
• CSRF защита
• Валидация данных
• Безопасная загрузка файлов"""
    
    content12_frame.text = tech_text
    for para in content12_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 13: Структура проекта ==========
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg13 = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg13.fill.solid()
    bg13.fill.fore_color.rgb = DARK_BG
    bg13.line.color.rgb = DARK_BG
    
    title13 = slide13.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title13_frame = title13.text_frame
    title13_frame.text = "СТРУКТУРА ПРОЕКТА"
    title13_frame.paragraphs[0].font.size = Pt(44)
    title13_frame.paragraphs[0].font.bold = True
    title13_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content13 = slide13.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content13_frame = content13.text_frame
    content13_frame.word_wrap = True
    
    structure_text = """📁 ОРГАНИЗАЦИЯ КОДА

uzmat/
├── uzmat_site/          # Главный проект Django
├── uzmat/               # Приложение
│   ├── models.py        # Модели данных
│   ├── views.py         # Логика представлений
│   ├── urls.py          # Маршрутизация
│   └── admin.py         # Админ-панель
├── templates/uzmat/     # HTML шаблоны (20+ страниц)
├── static/uzmat/        # CSS, JS, изображения
└── media/               # Загруженные файлы

📄 ОСНОВНЫЕ ФАЙЛЫ
• 20+ HTML шаблонов
• Единый стиль (styles.css)
• JavaScript для интерактивности
• Модели для всех сущностей"""
    
    content13_frame.text = structure_text
    for para in content13_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 14: Особенности дизайна ==========
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg14 = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg14.fill.solid()
    bg14.fill.fore_color.rgb = DARK_BG
    bg14.line.color.rgb = DARK_BG
    
    title14 = slide14.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title14_frame = title14.text_frame
    title14_frame.text = "ДИЗАЙН И UX"
    title14_frame.paragraphs[0].font.size = Pt(44)
    title14_frame.paragraphs[0].font.bold = True
    title14_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content14 = slide14.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content14_frame = content14.text_frame
    content14_frame.word_wrap = True
    
    design_text = """🎨 ЦВЕТОВАЯ СХЕМА
• Темный фон (премиум-стиль)
• Зеленые акценты (#009688)
• Белый текст для контраста
• Плавные переходы цветов

📱 АДАПТИВНОСТЬ
• Мобильные устройства (320px+)
• Планшеты (768px+)
• Десктоп (1024px+)
• Большие экраны (1440px+)

✨ АНИМАЦИИ
• Плавные переходы
• Hover-эффекты
• Загрузка контента
• Интерактивные элементы

🎯 UX ПРИНЦИПЫ
• Интуитивная навигация
• Быстрый доступ к функциям
• Минималистичный дизайн
• Фокус на контенте"""
    
    content14_frame.text = design_text
    for para in content14_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 15: Финальный ==========
    slide15 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg15 = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg15.fill.solid()
    bg15.fill.fore_color.rgb = DARK_BG
    bg15.line.color.rgb = DARK_BG
    
    # Заголовок
    final_title = slide15.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    final_title_frame = final_title.text_frame
    final_title_frame.text = "СПАСИБО ЗА ВНИМАНИЕ!"
    final_title_frame.paragraphs[0].font.size = Pt(48)
    final_title_frame.paragraphs[0].font.bold = True
    final_title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    final_title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Подзаголовок
    final_subtitle = slide15.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    final_subtitle_frame = final_subtitle.text_frame
    final_subtitle_frame.text = "UZMAT"
    final_subtitle_frame.paragraphs[0].font.size = Pt(36)
    final_subtitle_frame.paragraphs[0].font.color.rgb = LIGHT_TEXT
    final_subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Контакты
    final_contact = slide15.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    final_contact_frame = final_contact.text_frame
    final_contact_frame.text = "Премиальный маркетплейс спецтехники\nв Центральной Азии"
    final_contact_frame.paragraphs[0].font.size = Pt(20)
    final_contact_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    final_contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    final_contact_frame.paragraphs[1].font.size = Pt(20)
    final_contact_frame.paragraphs[1].font.color.rgb = RGBColor(150, 150, 150)
    final_contact_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Сохраняем презентацию
    filename = f"Uzmat_Presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(filename)
    print(f"✅ Презентация успешно создана: {filename}")
    print(f"📊 Всего слайдов: {len(prs.slides)}")
    return filename

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Ошибка: Не установлена библиотека python-pptx")
        print("📦 Установите её командой: pip install python-pptx")
    except Exception as e:
        print(f"❌ Ошибка при создании презентации: {e}")










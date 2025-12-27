#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Скрипт для создания презентации PowerPoint о проекте Uzmat с изображениями
Требует установки: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import os
from pathlib import Path

def add_image_to_slide(slide, image_path, left, top, width, height):
    """Добавляет изображение на слайд, если файл существует"""
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, left, top, width, height)
            return True
        except Exception as e:
            print(f"⚠️  Не удалось добавить изображение {image_path}: {e}")
            return False
    else:
        print(f"⚠️  Изображение не найдено: {image_path}")
        return False

def create_presentation():
    """Создает премиальную презентацию о проекте Uzmat с изображениями"""
    
    # Путь к папке со скриншотами
    base_dir = Path(__file__).parent
    screenshots_dir = base_dir / "screenshots"
    
    # Создаем папку для скриншотов, если её нет
    screenshots_dir.mkdir(exist_ok=True)
    
    # Создаем презентацию
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Цветовая схема (премиальная темная тема)
    PRIMARY_COLOR = RGBColor(0, 150, 136)  # Зеленый (как в проекте)
    DARK_BG = RGBColor(15, 23, 42)  # Темно-синий фон
    LIGHT_TEXT = RGBColor(255, 255, 255)  # Белый текст
    ACCENT_COLOR = RGBColor(34, 197, 94)  # Яркий зеленый акцент
    
    # ========== СЛАЙД 1: Титульный ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой слайд
    
    # Фон
    background = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.color.rgb = DARK_BG
    
    # Логотип (если есть)
    logo_path = base_dir / "uzmat" / "logo-uzmat.svg"
    if not logo_path.exists():
        logo_path = base_dir / "logo-uzmat.svg"
    
    # Заголовок
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "UZMAT"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Подзаголовок
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Премиальный маркетплейс спецтехники"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = LIGHT_TEXT
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Дата
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime("%d %B %Y")
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(16)
    date_para.font.color.rgb = RGBColor(150, 150, 150)
    date_para.alignment = PP_ALIGN.CENTER
    
    # ========== СЛАЙД 2: О проекте ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = DARK_BG
    bg2.line.color.rgb = DARK_BG
    
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title2_frame = title2.text_frame
    title2_frame.text = "О ПРОЕКТЕ"
    title2_frame.paragraphs[0].font.size = Pt(44)
    title2_frame.paragraphs[0].font.bold = True
    title2_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Описание слева
    desc_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    
    description = """Uzmat — это премиальная платформа для поиска и размещения объявлений о спецтехнике в Центральной Азии.

🎯 ОСНОВНАЯ ЦЕЛЬ
Единая точка доступа для поиска надежных подрядчиков и управления парком техники

✨ ОСОБЕННОСТИ
• Современный премиальный дизайн
• Полная адаптивность
• Плавные анимации
• Многоязычная поддержка
• Безопасная система

🌍 ГЕОГРАФИЯ
Казахстан • Узбекистан • Россия"""
    
    desc_frame.text = description
    for para in desc_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(12)
    
    # Скриншот главной страницы справа
    main_screenshot = screenshots_dir / "main_page.png"
    if not main_screenshot.exists():
        main_screenshot = screenshots_dir / "main_page.jpg"
    if main_screenshot.exists():
        add_image_to_slide(slide2, str(main_screenshot), Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    
    # ========== СЛАЙД 3: Основные функции ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = DARK_BG
    bg3.line.color.rgb = DARK_BG
    
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title3_frame = title3.text_frame
    title3_frame.text = "ОСНОВНЫЕ ФУНКЦИИ"
    title3_frame.paragraphs[0].font.size = Pt(44)
    title3_frame.paragraphs[0].font.bold = True
    title3_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Функции в две колонки
    functions_left = [
        "🏗️ Продажа спецтехники",
        "🚛 Аренда спецтехники",
        "🔧 Запчасти",
        "🛠️ Ремонт и услуги"
    ]
    
    functions_right = [
        "💬 Система чатов",
        "⭐ Избранное",
        "👤 Профиль пользователя",
        "🔍 Умный поиск и фильтры"
    ]
    
    # Левая колонка
    left_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.text = "\n".join(functions_left)
    for para in left_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(16)
        para.font.bold = True
    
    # Правая колонка
    right_box = slide3.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.text = "\n".join(functions_right)
    for para in right_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(16)
        para.font.bold = True
    
    # ========== СЛАЙД 4: Главная страница ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = DARK_BG
    bg4.line.color.rgb = DARK_BG
    
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title4_frame = title4.text_frame
    title4_frame.text = "ГЛАВНАЯ СТРАНИЦА"
    title4_frame.paragraphs[0].font.size = Pt(44)
    title4_frame.paragraphs[0].font.bold = True
    title4_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Скриншот главной страницы (большой)
    main_full = screenshots_dir / "index_full.png"
    if not main_full.exists():
        main_full = screenshots_dir / "index_full.jpg"
    if main_full.exists():
        add_image_to_slide(slide4, str(main_full), Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
    else:
        # Если нет скриншота, добавляем текст
        content4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content4_frame = content4.text_frame
        content4_frame.word_wrap = True
        
        main_page_text = """📍 URL: / (главная страница)

🎨 ДИЗАЙН
• Премиальный темный интерфейс
• Hero-секция с поиском
• Быстрые кнопки популярных запросов

📋 СЕКЦИИ
1. Горячие предложения
2. Популярные предложения
3. Категории (Продажа, Аренда, Заявки, Запчасти, Ремонт)

🔍 ФУНКЦИОНАЛ
• Поиск в реальном времени
• Фильтры по стране, городу, типу, марке, цене
• AJAX-фильтрация без перезагрузки"""
        
        content4_frame.text = main_page_text
        for para in content4_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.color.rgb = LIGHT_TEXT
            para.space_after = Pt(10)
    
    # ========== СЛАЙД 5: Каталог ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = DARK_BG
    bg5.line.color.rgb = DARK_BG
    
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title5_frame = title5.text_frame
    title5_frame.text = "КАТАЛОГ ОБЪЯВЛЕНИЙ"
    title5_frame.paragraphs[0].font.size = Pt(44)
    title5_frame.paragraphs[0].font.bold = True
    title5_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Скриншот каталога
    catalog_img = screenshots_dir / "catalog.png"
    if not catalog_img.exists():
        catalog_img = screenshots_dir / "catalog.jpg"
    if catalog_img.exists():
        add_image_to_slide(slide5, str(catalog_img), Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
    else:
        content5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content5_frame = content5.text_frame
        content5_frame.word_wrap = True
        
        catalog_text = """📍 URL: /catalog/

🎯 НАЗНАЧЕНИЕ
Централизованный каталог всех объявлений о продаже и аренде

📊 РЕЖИМЫ
• Все объявления (продажа + аренда)
• Только продажа - /catalog/?ad_type=sale
• Только аренда - /catalog/?ad_type=rent

🔍 ФИЛЬТРЫ
✓ Страна, Город, Тип, Марка, Цена
✓ AJAX-фильтрация без перезагрузки
✓ Пагинация (12 объявлений на страницу)"""
        
        content5_frame.text = catalog_text
        for para in content5_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.color.rgb = LIGHT_TEXT
            para.space_after = Pt(10)
    
    # ========== СЛАЙД 6: Запчасти ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = DARK_BG
    bg6.line.color.rgb = DARK_BG
    
    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title6_frame = title6.text_frame
    title6_frame.text = "ЗАПЧАСТИ"
    title6_frame.paragraphs[0].font.size = Pt(44)
    title6_frame.paragraphs[0].font.bold = True
    title6_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Скриншот страницы запчастей
    parts_img = screenshots_dir / "parts.png"
    if not parts_img.exists():
        parts_img = screenshots_dir / "parts.jpg"
    if parts_img.exists():
        add_image_to_slide(slide6, str(parts_img), Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
    else:
        content6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content6_frame = content6.text_frame
        content6_frame.word_wrap = True
        
        parts_text = """📍 URL: /parts/

🎯 НАЗНАЧЕНИЕ
Специализированный раздел для поиска запчастей

🔍 ПОИСК
• Поиск по названию
• Быстрые кнопки популярных запросов
• Фильтры по стране, городу, типу, марке, цене

📊 ОТОБРАЖЕНИЕ
• Сетка карточек объявлений
• Пагинация
• Динамическая фильтрация"""
        
        content6_frame.text = parts_text
        for para in content6_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.color.rgb = LIGHT_TEXT
            para.space_after = Pt(10)
    
    # ========== СЛАЙД 7: Ремонт ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg7 = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg7.fill.solid()
    bg7.fill.fore_color.rgb = DARK_BG
    bg7.line.color.rgb = DARK_BG
    
    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title7_frame = title7.text_frame
    title7_frame.text = "РЕМОНТ И УСЛУГИ"
    title7_frame.paragraphs[0].font.size = Pt(44)
    title7_frame.paragraphs[0].font.bold = True
    title7_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Скриншот страницы ремонта
    repair_img = screenshots_dir / "repair.png"
    if not repair_img.exists():
        repair_img = screenshots_dir / "repair.jpg"
    if repair_img.exists():
        add_image_to_slide(slide7, str(repair_img), Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
    else:
        content7 = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content7_frame = content7.text_frame
        content7_frame.word_wrap = True
        
        repair_text = """📍 URL: /logistics/

🎯 НАЗНАЧЕНИЕ
Раздел для поиска услуг по ремонту спецтехники

📋 КОНТЕНТ
• Все объявления типа "Услуги"
• Специализированные сервисы
• Профессиональные мастера

🔍 ФИЛЬТРЫ
✓ Страна, Город, Тип, Марка, Цена
✓ Динамическая фильтрация
✓ Пагинация результатов"""
        
        content7_frame.text = repair_text
        for para in content7_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.color.rgb = LIGHT_TEXT
            para.space_after = Pt(10)
    
    # ========== СЛАЙД 8: Профиль ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg8.fill.solid()
    bg8.fill.fore_color.rgb = DARK_BG
    bg8.line.color.rgb = DARK_BG
    
    title8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title8_frame = title8.text_frame
    title8_frame.text = "ПРОФИЛЬ ПОЛЬЗОВАТЕЛЯ"
    title8_frame.paragraphs[0].font.size = Pt(44)
    title8_frame.paragraphs[0].font.bold = True
    title8_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Скриншот профиля
    profile_img = screenshots_dir / "profile.png"
    if not profile_img.exists():
        profile_img = screenshots_dir / "profile.jpg"
    if profile_img.exists():
        add_image_to_slide(slide8, str(profile_img), Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
    else:
        content8 = slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content8_frame = content8.text_frame
        content8_frame.word_wrap = True
        
        auth_text = """📍 URL: /profile/

👤 ФУНКЦИИ
✓ Просмотр своих объявлений
✓ Управление объявлениями
✓ Редактирование профиля
✓ Настройки аккаунта
✓ История активности

🔐 АВТОРИЗАЦИЯ
• Регистрация (физ. лицо/компания)
• Безопасный вход
• Восстановление пароля"""
        
        content8_frame.text = auth_text
        for para in content8_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.color.rgb = LIGHT_TEXT
            para.space_after = Pt(10)
    
    # ========== СЛАЙД 9: Создание объявления ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg9.fill.solid()
    bg9.fill.fore_color.rgb = DARK_BG
    bg9.line.color.rgb = DARK_BG
    
    title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title9_frame = title9.text_frame
    title9_frame.text = "СОЗДАНИЕ ОБЪЯВЛЕНИЯ"
    title9_frame.paragraphs[0].font.size = Pt(44)
    title9_frame.paragraphs[0].font.bold = True
    title9_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Скриншот формы создания
    create_img = screenshots_dir / "create_ad.png"
    if not create_img.exists():
        create_img = screenshots_dir / "create_ad.jpg"
    if create_img.exists():
        add_image_to_slide(slide9, str(create_img), Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
    else:
        content9 = slide9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content9_frame = content9.text_frame
        content9_frame.word_wrap = True
        
        create_text = """📍 URL: /create/

📝 ТИПЫ ОБЪЯВЛЕНИЙ
• Продажа • Аренда • Запчасти • Услуги

📋 ФОРМА
✓ Заголовок и описание
✓ Тип техники и марка
✓ Год выпуска и состояние
✓ Цена и валюта
✓ Местоположение
✓ Контакты
✓ Загрузка фотографий

⚡ ОСОБЕННОСТИ
• Валидация данных
• Автоматический slug
• Предпросмотр"""
        
        content9_frame.text = create_text
        for para in content9_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.color.rgb = LIGHT_TEXT
            para.space_after = Pt(10)
    
    # ========== СЛАЙД 10: Детальная страница ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg10.fill.solid()
    bg10.fill.fore_color.rgb = DARK_BG
    bg10.line.color.rgb = DARK_BG
    
    title10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title10_frame = title10.text_frame
    title10_frame.text = "ДЕТАЛЬНАЯ СТРАНИЦА"
    title10_frame.paragraphs[0].font.size = Pt(44)
    title10_frame.paragraphs[0].font.bold = True
    title10_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Скриншот детальной страницы
    detail_img = screenshots_dir / "ad_detail.png"
    if not detail_img.exists():
        detail_img = screenshots_dir / "ad_detail.jpg"
    if detail_img.exists():
        add_image_to_slide(slide10, str(detail_img), Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
    else:
        content10 = slide10.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content10_frame = content10.text_frame
        content10_frame.word_wrap = True
        
        detail_text = """📍 URL: /ad/<slug>/

📸 ГАЛЕРЕЯ
• Множественные фотографии
• Увеличение по клику

📋 ИНФОРМАЦИЯ
✓ Полное описание
✓ Характеристики
✓ Цена и условия
✓ Местоположение
✓ Контакты

⭐ ФУНКЦИИ
• Избранное
• Поделиться
• Связаться с продавцом
• Другие объявления пользователя"""
        
        content10_frame.text = detail_text
        for para in content10_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.color.rgb = LIGHT_TEXT
            para.space_after = Pt(10)
    
    # ========== СЛАЙД 11: Чаты ==========
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg11 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg11.fill.solid()
    bg11.fill.fore_color.rgb = DARK_BG
    bg11.line.color.rgb = DARK_BG
    
    title11 = slide11.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title11_frame = title11.text_frame
    title11_frame.text = "ЧАТЫ И ЗАЯВКИ"
    title11_frame.paragraphs[0].font.size = Pt(44)
    title11_frame.paragraphs[0].font.bold = True
    title11_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Скриншот чатов
    chats_img = screenshots_dir / "chats.png"
    if not chats_img.exists():
        chats_img = screenshots_dir / "chats.jpg"
    if chats_img.exists():
        add_image_to_slide(slide11, str(chats_img), Inches(0.5), Inches(1.3), Inches(9), Inches(5.7))
    else:
        content11 = slide11.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content11_frame = content11.text_frame
        content11_frame.word_wrap = True
        
        chats_text = """📍 URL: /chats/

💬 ФУНКЦИОНАЛ
• Прямая связь покупатель-продавец
• Обмен сообщениями
• История переписки

📋 ЗАЯВКИ
• Создание заявок на аренду
• Управление заявками
• Статусы заявок

🔔 УВЕДОМЛЕНИЯ
• Новые сообщения
• Обновления заявок"""
        
        content11_frame.text = chats_text
        for para in content11_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.color.rgb = LIGHT_TEXT
            para.space_after = Pt(10)
    
    # ========== СЛАЙД 12: Технологии ==========
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg12.fill.solid()
    bg12.fill.fore_color.rgb = DARK_BG
    bg12.line.color.rgb = DARK_BG
    
    title12 = slide12.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title12_frame = title12.text_frame
    title12_frame.text = "ТЕХНОЛОГИИ"
    title12_frame.paragraphs[0].font.size = Pt(44)
    title12_frame.paragraphs[0].font.bold = True
    title12_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content12 = slide12.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content12_frame = content12.text_frame
    content12_frame.word_wrap = True
    
    tech_text = """🖥️ BACKEND
• Django 4.2.7
• Python 3.12
• SQLite
• Django REST Framework

🎨 FRONTEND
• HTML5, CSS3, JavaScript
• Адаптивный дизайн
• AJAX для динамической загрузки
• Плавные анимации

🗄️ БАЗА ДАННЫХ
• Модели: User, Advertisement, Category, Favorite
• Связи и индексы
• Миграции Django

🔒 БЕЗОПАСНОСТЬ
• Аутентификация Django
• CSRF защита
• Валидация данных
• Безопасная загрузка файлов"""
    
    content12_frame.text = tech_text
    for para in content12_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 13: Дизайн ==========
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg13 = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg13.fill.solid()
    bg13.fill.fore_color.rgb = DARK_BG
    bg13.line.color.rgb = DARK_BG
    
    title13 = slide13.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title13_frame = title13.text_frame
    title13_frame.text = "ДИЗАЙН И UX"
    title13_frame.paragraphs[0].font.size = Pt(44)
    title13_frame.paragraphs[0].font.bold = True
    title13_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content13 = slide13.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content13_frame = content13.text_frame
    content13_frame.word_wrap = True
    
    design_text = """🎨 ЦВЕТОВАЯ СХЕМА
• Темный фон (премиум-стиль)
• Зеленые акценты (#009688)
• Белый текст для контраста
• Плавные переходы

📱 АДАПТИВНОСТЬ
• Мобильные (320px+)
• Планшеты (768px+)
• Десктоп (1024px+)
• Большие экраны (1440px+)

✨ АНИМАЦИИ
• Плавные переходы
• Hover-эффекты
• Загрузка контента
• Интерактивные элементы

🎯 UX ПРИНЦИПЫ
• Интуитивная навигация
• Быстрый доступ
• Минималистичный дизайн
• Фокус на контенте"""
    
    content13_frame.text = design_text
    for para in content13_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = LIGHT_TEXT
        para.space_after = Pt(10)
    
    # ========== СЛАЙД 14: Финальный ==========
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg14 = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg14.fill.solid()
    bg14.fill.fore_color.rgb = DARK_BG
    bg14.line.color.rgb = DARK_BG
    
    final_title = slide14.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    final_title_frame = final_title.text_frame
    final_title_frame.text = "СПАСИБО ЗА ВНИМАНИЕ!"
    final_title_frame.paragraphs[0].font.size = Pt(48)
    final_title_frame.paragraphs[0].font.bold = True
    final_title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    final_title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    final_subtitle = slide14.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    final_subtitle_frame = final_subtitle.text_frame
    final_subtitle_frame.text = "UZMAT"
    final_subtitle_frame.paragraphs[0].font.size = Pt(36)
    final_subtitle_frame.paragraphs[0].font.color.rgb = LIGHT_TEXT
    final_subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    final_contact = slide14.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    final_contact_frame = final_contact.text_frame
    final_contact_frame.text = "Премиальный маркетплейс спецтехники\nв Центральной Азии"
    final_contact_frame.paragraphs[0].font.size = Pt(20)
    final_contact_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    final_contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    final_contact_frame.paragraphs[1].font.size = Pt(20)
    final_contact_frame.paragraphs[1].font.color.rgb = RGBColor(150, 150, 150)
    final_contact_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # Сохраняем презентацию
    filename = f"Uzmat_Presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(filename)
    print(f"✅ Презентация успешно создана: {filename}")
    print(f"📊 Всего слайдов: {len(prs.slides)}")
    print(f"\n📸 ИНСТРУКЦИЯ ПО СКРИНШОТАМ:")
    print(f"   Поместите скриншоты в папку: {screenshots_dir}")
    print(f"   Имена файлов:")
    print(f"   - main_page.png/jpg - главная страница")
    print(f"   - index_full.png/jpg - полный скриншот главной")
    print(f"   - catalog.png/jpg - каталог")
    print(f"   - parts.png/jpg - запчасти")
    print(f"   - repair.png/jpg - ремонт")
    print(f"   - profile.png/jpg - профиль")
    print(f"   - create_ad.png/jpg - создание объявления")
    print(f"   - ad_detail.png/jpg - детальная страница")
    print(f"   - chats.png/jpg - чаты")
    print(f"\n   Если скриншоты не найдены, будут использованы текстовые описания.")
    return filename

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Ошибка: Не установлена библиотека python-pptx")
        print("📦 Установите её командой: pip install python-pptx")
    except Exception as e:
        print(f"❌ Ошибка при создании презентации: {e}")
        import traceback
        traceback.print_exc()









